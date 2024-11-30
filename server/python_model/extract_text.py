import argparse
import os
from pptx import Presentation
from PyPDF2 import PdfReader
import docx
import pandas as pd
import json

def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = [page.extract_text() for page in reader.pages]
        return "\n".join(text)
    except Exception as e:
        return f"Error extracting text from PDF: {e}"

def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Error extracting text from DOCX: {e}"

def extract_text_from_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    except Exception as e:
        return f"Error extracting text from CSV: {e}"

def extract_text_from_excel(file_path):
    try:
        df = pd.read_excel(file_path)
        return df.to_string(index=False)
    except Exception as e:
        return f"Error extracting text from Excel: {e}"

def extract_text_from_json(file_path):
    try:
        with open(file_path, 'r') as f:
            data = json.load(f)
        return json.dumps(data, indent=4)
    except Exception as e:
        return f"Error extracting text from JSON: {e}"

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error extracting text from PPTX: {e}"

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract text from various file formats")
    parser.add_argument("--pdf", type=str, help="Path to the PDF file")
    parser.add_argument("--docx", type=str, help="Path to the DOCX file")
    parser.add_argument("--csv", type=str, help="Path to the CSV file")
    parser.add_argument("--excel", type=str, help="Path to the Excel file")
    parser.add_argument("--json", type=str, help="Path to the JSON file")
    parser.add_argument("--pptx", type=str, help="Path to the PPTX file")
    parser.add_argument("--txt", type=str, help="Path to the TXT file")
    args = parser.parse_args()

    if args.pdf:
        print(extract_text_from_pdf(args.pdf))
    elif args.docx:
        print(extract_text_from_docx(args.docx))
    elif args.csv:
        print(extract_text_from_csv(args.csv))
    elif args.excel:
        print(extract_text_from_excel(args.excel))
    elif args.json:
        print(extract_text_from_json(args.json))
    elif args.pptx:
        print(extract_text_from_pptx(args.pptx))
    elif args.txt:
        try:
            with open(args.txt, 'r') as file:
                print(file.read())
        except Exception as e:
            print(f"Error reading text file: {e}")
    else:
        print("No valid file format specified.")
