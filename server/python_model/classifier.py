import sys
import os
import joblib
import numpy as np
from preprocess import clean_text  
from PyPDF2 import PdfReader
from pptx import Presentation
import docx  
import pandas as pd  
import json
import nltk

script_dir = os.path.dirname(os.path.abspath(__file__))

model_filename = os.path.join(script_dir, 'model.pkl')
vectorizer_filename = os.path.join(script_dir, 'vectorizer.pkl')

def extract_text_from_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""  
        return text.strip()
    except Exception as e:
        raise ValueError(f"Unable to extract text from {pdf_path}. Error: {e}")

def extract_text_from_pptx(pptx_path):
    try:
        presentation = Presentation(pptx_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += " ".join(run.text for run in paragraph.runs) + " "
        return text.strip()
    except Exception as e:
        raise ValueError(f"Unable to extract text from {pptx_path}. Error: {e}")

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        text = " ".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise ValueError(f"Unable to extract text from {docx_path}. Error: {e}")

def extract_text_from_xlsx(xlsx_path):
    try:
        df = pd.read_excel(xlsx_path)
        text = " ".join(df.to_string(index=False))
        return text.strip()
    except Exception as e:
        raise ValueError(f"Unable to extract text from {xlsx_path}. Error: {e}")

def extract_text_from_json(json_path):
    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            text = str(data)
        return text.strip()
    except Exception as e:
        raise ValueError(f"Unable to extract text from {json_path}. Error: {e}")

def classify_document(file_path):
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File {file_path} does not exist.")

    # Load the model and vectorizer
    try:
        model = joblib.load(model_filename)
        vectorizer = joblib.load(vectorizer_filename)
    except FileNotFoundError as e:
        raise FileNotFoundError(f"Model or vectorizer file not found: {e}")

    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.pdf':
            content = extract_text_from_pdf(file_path)
        elif ext == '.pptx':
            content = extract_text_from_pptx(file_path)
        elif ext == '.docx':
            content = extract_text_from_docx(file_path)
        elif ext == '.xlsx':
            content = extract_text_from_xlsx(file_path)
        elif ext == '.json':
            content = extract_text_from_json(file_path)
        elif ext in ['.txt', '.csv']:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return {"category": "Unclassified", "confidence": 0.0, "probabilities": {}}

    if not content or len(content.strip()) == 0:
        print(f"Warning: No extractable text from {file_path}")
        return {"category": "Unclassified", "confidence": 0.0, "probabilities": {}}

    try:
        cleaned_content = clean_text(content)
    except Exception as e:
        print(f"Error during text preprocessing for {file_path}: {e}")
        return {"category": "Unclassified", "confidence": 0.0, "probabilities": {}}

    # Vectorize and classify with confidence
    try:
        vectorized_content = vectorizer.transform([cleaned_content])
        
        if hasattr(model, 'predict_proba'):
            probabilities = model.predict_proba(vectorized_content)[0]
            prediction = model.classes_[np.argmax(probabilities)]
            confidence = float(np.max(probabilities))
            
            prob_dict = dict(zip(model.classes_, probabilities.tolist()))
        else:
            prediction = model.predict(vectorized_content)[0]
            confidence = 1.0  
            prob_dict = {}

        return {
            "category": prediction, 
            "confidence": confidence, 
            "probabilities": prob_dict
        }
    except Exception as e:
        print(f"Error during classification for {file_path}: {e}")
        return {"category": "Unclassified", "confidence": 0.0, "probabilities": {}}

def classify_directory(directory_path):
    """
    Classify all supported documents in a directory (including subdirectories).
    """
    supported_extensions = ['.pdf', '.pptx', '.docx', '.xlsx', '.json', '.txt', '.csv']
    
    classification_results = {}
    
    for root, dirs, files in os.walk(directory_path):
        for file in files:
            file_path = os.path.join(root, file)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext in supported_extensions:
                try:
                    classification = classify_document(file_path)
                    classification_results[file_path] = classification
                except Exception as e:
                    print(f"Failed to classify {file_path}: {e}")
    
    return classification_results

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python classifier.py <file_or_directory_to_classify>")
        sys.exit(1)

    input_path = sys.argv[1]
    
    try:
        # Check if it's a directory or a single file
        if os.path.isdir(input_path):
            # Classify all documents in the directory
            results = classify_directory(input_path)
            
            print("Classification Results:")
            for file_path, result in results.items():
                print(f"{file_path}: Category={result['category']}, Confidence={result['confidence']:.4f}")
        else:
            # Classify a single file
            result = classify_document(input_path)
            print(f"Category: {result['category']}")
            print(f"Confidence: {result['confidence']:.4f}")
            print("Probabilities:", result['probabilities'])
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)